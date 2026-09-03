"""Extract searchable text from supported document formats."""

from io import BytesIO
from pathlib import Path
import shutil
import subprocess
import tempfile

from docx import Document
from pypdf import PdfReader
from pptx import Presentation


SUPPORTED_EXTENSIONS = {".txt", ".md", ".json", ".pdf", ".doc", ".docx", ".ppt", ".pptx"}


def extract_text(contents: bytes, filename: str) -> str:
    suffix = Path(filename).suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        raise ValueError("Upload a text, PDF, Word, or PowerPoint file (.txt, .md, .json, .pdf, .doc, .docx, .ppt, .pptx)")

    if suffix in {".txt", ".md", ".json"}:
        return contents.decode("utf-8")
    if suffix == ".pdf":
        return "\n\n".join(page.extract_text() or "" for page in PdfReader(BytesIO(contents)).pages).strip()
    if suffix == ".docx":
        document = Document(BytesIO(contents))
        return "\n\n".join(paragraph.text for paragraph in document.paragraphs).strip()
    if suffix == ".pptx":
        presentation = Presentation(BytesIO(contents))
        slides = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            text = "\n".join(
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            )
            if text:
                slides.append(f"Slide {slide_number}\n{text}")
        return "\n\n".join(slides).strip()

    converter = shutil.which("antiword") if suffix == ".doc" else shutil.which("soffice") or shutil.which("libreoffice")
    if converter is None:
        required_tool = "antiword" if suffix == ".doc" else "LibreOffice"
        raise RuntimeError(f"Legacy {suffix} files require {required_tool}; install it or convert the file to {suffix}x")
    if suffix == ".doc":
        result = subprocess.run(
            [converter, "-"],
            input=contents,
            capture_output=True,
            check=False,
        )
        if result.returncode != 0:
            raise ValueError("Could not read the legacy .doc file")
        return result.stdout.decode("utf-8", errors="replace").strip()

    with tempfile.TemporaryDirectory() as directory:
        source = Path(directory) / Path(filename).name
        source.write_bytes(contents)
        result = subprocess.run(
            [converter, "--headless", "--convert-to", "pptx", "--outdir", directory, str(source)],
            capture_output=True,
            check=False,
        )
        converted = source.with_suffix(".pptx")
        if result.returncode != 0 or not converted.exists():
            raise ValueError("Could not read the legacy .ppt file")
        return extract_text(converted.read_bytes(), converted.name)
