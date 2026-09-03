from io import BytesIO

from pptx import Presentation

from app.ingestion.document_text import SUPPORTED_EXTENSIONS, extract_text


def test_extract_text_from_pptx():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Release Plan"
    textbox = slide.shapes.add_textbox(0, 0, 100, 100)
    textbox.text = "Deployment checklist"

    contents = BytesIO()
    presentation.save(contents)

    text = extract_text(contents.getvalue(), "release-plan.pptx")

    assert "Slide 1" in text
    assert "Release Plan" in text
    assert "Deployment checklist" in text


def test_pptx_is_supported_extension():
    assert ".pptx" in SUPPORTED_EXTENSIONS
    assert ".ppt" in SUPPORTED_EXTENSIONS
